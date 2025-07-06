# pdftools/utils.py
import os
import tempfile
import fitz  # PyMuPDF
from PIL import Image
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptxInches
from pptx.enum.text import PP_ALIGN
from PyPDF2 import PdfWriter, PdfReader
import magic
import io
import logging

logger = logging.getLogger(__name__)

def validate_pdf_file(file):
    """Validate if uploaded file is a valid PDF"""
    try:
        # Check file extension
        if not file.name.lower().endswith('.pdf'):
            return False
        
        # Check file size (max 50MB)
        if file.size > 50 * 1024 * 1024:
            return False
        
        # Check MIME type
        file_type = magic.from_buffer(file.read(1024), mime=True)
        file.seek(0)  # Reset file pointer
        
        return file_type == 'application/pdf'
    
    except Exception as e:
        logger.error(f"PDF validation error: {str(e)}")
        return False

def merge_pdfs(input_files, task_id):
    """Merge multiple PDF files into one"""
    try:
        output_path = os.path.join(tempfile.gettempdir(), f'merged_{task_id}.pdf')
        
        pdf_writer = PdfWriter()
        
        for file_path in input_files:
            with open(file_path, 'rb') as file:
                pdf_reader = PdfReader(file)
                
                # Add all pages from current PDF
                for page_num in range(len(pdf_reader.pages)):
                    page = pdf_reader.pages[page_num]
                    pdf_writer.add_page(page)
        
        # Write merged PDF
        with open(output_path, 'wb') as output_file:
            pdf_writer.write(output_file)
        
        return output_path
    
    except Exception as e:
        logger.error(f"PDF merge error: {str(e)}")
        raise Exception(f"Failed to merge PDFs: {str(e)}")

def convert_pdf_to_word(pdf_path, task_id):
    """Convert PDF to Word document preserving images and formatting"""
    try:
        output_path = os.path.join(tempfile.gettempdir(), f'converted_{task_id}.docx')
        
        # Open PDF document
        pdf_document = fitz.open(pdf_path)
        
        # Create Word document
        doc = Document()
        
        # Set document margins
        sections = doc.sections
        for section in sections:
            section.top_margin = Inches(0.5)
            section.bottom_margin = Inches(0.5)
            section.left_margin = Inches(0.5)
            section.right_margin = Inches(0.5)
        
        # Process each page
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            # Add page break if not first page
            if page_num > 0:
                doc.add_page_break()
            
            # Extract text blocks
            text_blocks = page.get_text("dict")
            
            # Extract images
            image_list = page.get_images(full=True)
            
            # Process text blocks
            for block in text_blocks["blocks"]:
                if "lines" in block:
                    # Text block
                    block_text = ""
                    for line in block["lines"]:
                        for span in line["spans"]:
                            block_text += span["text"]
                        block_text += "\n"
                    
                    if block_text.strip():
                        # Add paragraph with extracted text
                        paragraph = doc.add_paragraph(block_text.strip())
                        
                        # Try to preserve some formatting
                        if block["lines"]:
                            first_span = block["lines"][0]["spans"][0]
                            font_size = first_span.get("size", 12)
                            
                            # Set font size (convert to points)
                            for run in paragraph.runs:
                                run.font.size = Pt(font_size)
            
            # Process images
            for img_index, img in enumerate(image_list):
                try:
                    # Get image data
                    xref = img[0]
                    pix = fitz.Pixmap(pdf_document, xref)
                    
                    # Convert to PIL Image
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("png")
                        
                        # Save image temporarily
                        temp_img_path = os.path.join(
                            tempfile.gettempdir(), 
                            f'temp_img_{task_id}_{page_num}_{img_index}.png'
                        )
                        
                        with open(temp_img_path, 'wb') as img_file:
                            img_file.write(img_data)
                        
                        # Add image to document
                        try:
                            # Calculate image size (max width 6 inches)
                            img_pil = Image.open(temp_img_path)
                            img_width, img_height = img_pil.size
                            
                            # Calculate scaling
                            max_width = 6.0  # inches
                            if img_width > max_width * 72:  # 72 DPI
                                scale = (max_width * 72) / img_width
                                doc_width = max_width
                                doc_height = (img_height * scale) / 72
                            else:
                                doc_width = img_width / 72
                                doc_height = img_height / 72
                            
                            # Add image to document
                            paragraph = doc.add_paragraph()
                            paragraph.alignment = WD_ALIGN_PARAGRAPH.CENTER
                            run = paragraph.runs[0] if paragraph.runs else paragraph.add_run()
                            run.add_picture(temp_img_path, width=Inches(doc_width))
                            
                            # Clean up temporary image
                            os.remove(temp_img_path)
                            
                        except Exception as img_error:
                            logger.warning(f"Failed to add image {img_index}: {str(img_error)}")
                            if os.path.exists(temp_img_path):
                                os.remove(temp_img_path)
                    
                    pix = None
                    
                except Exception as img_error:
                    logger.warning(f"Failed to process image {img_index}: {str(img_error)}")
                    continue
        
        # Save document
        doc.save(output_path)
        pdf_document.close()
        
        return output_path
    
    except Exception as e:
        logger.error(f"PDF to Word conversion error: {str(e)}")
        raise Exception(f"Failed to convert PDF to Word: {str(e)}")

def convert_pdf_to_ppt(pdf_path, task_id):
    """Convert PDF to PowerPoint presentation preserving images and formatting"""
    try:
        output_path = os.path.join(tempfile.gettempdir(), f'converted_{task_id}.pptx')
        
        # Open PDF document
        pdf_document = fitz.open(pdf_path)
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Process each page as a slide
        for page_num in range(len(pdf_document)):
            page = pdf_document[page_num]
            
            # Add slide with blank layout
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            
            # Get page dimensions
            page_rect = page.rect
            page_width = page_rect.width
            page_height = page_rect.height
            
            # Calculate slide dimensions (PowerPoint slide is 10" x 7.5")
            slide_width = PptxInches(10)
            slide_height = PptxInches(7.5)
            
            # Extract text blocks
            text_blocks = page.get_text("dict")
            
            # Extract images
            image_list = page.get_images(full=True)
            
            # Process text blocks
            text_content = []
            for block in text_blocks["blocks"]:
                if "lines" in block:
                    block_text = ""
                    for line in block["lines"]:
                        for span in line["spans"]:
                            block_text += span["text"]
                        block_text += "\n"
                    
                    if block_text.strip():
                        text_content.append(block_text.strip())
            
            # Add text to slide if any
            if text_content:
                # Create text box
                left = PptxInches(0.5)
                top = PptxInches(0.5)
                width = PptxInches(9)
                height = PptxInches(2)
                
                text_box = slide.shapes.add_textbox(left, top, width, height)
                text_frame = text_box.text_frame
                text_frame.word_wrap = True
                
                # Add all text content
                full_text = "\n\n".join(text_content)
                text_frame.text = full_text
                
                # Format text
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(14)
                    paragraph.alignment = PP_ALIGN.LEFT
            
            # Process images
            for img_index, img in enumerate(image_list):
                try:
                    # Get image data
                    xref = img[0]
                    pix = fitz.Pixmap(pdf_document, xref)
                    
                    # Convert to PIL Image
                    if pix.n - pix.alpha < 4:  # GRAY or RGB
                        img_data = pix.tobytes("png")
                        
                        # Save image temporarily
                        temp_img_path = os.path.join(
                            tempfile.gettempdir(), 
                            f'temp_img_{task_id}_{page_num}_{img_index}.png'
                        )
                        
                        with open(temp_img_path, 'wb') as img_file:
                            img_file.write(img_data)
                        
                        # Add image to slide
                        try:
                            # Calculate image position and size
                            img_pil = Image.open(temp_img_path)
                            img_width, img_height = img_pil.size
                            
                            # Calculate scaling to fit slide
                            max_width = 8.0  # inches
                            max_height = 5.0  # inches
                            
                            # Calculate aspect ratio preserving scale
                            width_scale = max_width / (img_width / 72)
                            height_scale = max_height / (img_height / 72)
                            scale = min(width_scale, height_scale, 1.0)
                            
                            final_width = (img_width / 72) * scale
                            final_height = (img_height / 72) * scale
                            
                            # Center image on slide
                            left = (10 - final_width) / 2
                            top = (7.5 - final_height) / 2
                            
                            # Add image to slide
                            slide.shapes.add_picture(
                                temp_img_path, 
                                PptxInches(left), 
                                PptxInches(top),
                                width=PptxInches(final_width),
                                height=PptxInches(final_height)
                            )
                            
                            # Clean up temporary image
                            os.remove(temp_img_path)
                            
                        except Exception as img_error:
                            logger.warning(f"Failed to add image {img_index} to slide: {str(img_error)}")
                            if os.path.exists(temp_img_path):
                                os.remove(temp_img_path)
                    
                    pix = None
                    
                except Exception as img_error:
                    logger.warning(f"Failed to process image {img_index} on page {page_num}: {str(img_error)}")
                    continue
        
        # Save presentation
        prs.save(output_path)
        pdf_document.close()
        
        return output_path
    
    except Exception as e:
        logger.error(f"PDF to PowerPoint conversion error: {str(e)}")
        raise Exception(f"Failed to convert PDF to PowerPoint: {str(e)}")

def password_protect_pdf(pdf_path, password, task_id):
    """Add password protection to PDF"""
    try:
        output_path = os.path.join(tempfile.gettempdir(), f'protected_{task_id}.pdf')
        
        # Read original PDF
        with open(pdf_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            pdf_writer = PdfWriter()
            
            # Add all pages
            for page_num in range(len(pdf_reader.pages)):
                pdf_writer.add_page(pdf_reader.pages[page_num])
            
            # Add password protection
            pdf_writer.encrypt(password)
            
            # Write protected PDF
            with open(output_path, 'wb') as output_file:
                pdf_writer.write(output_file)
        
        return output_path
    
    except Exception as e:
        logger.error(f"PDF password protection error: {str(e)}")
        raise Exception(f"Failed to password protect PDF: {str(e)}")

def extract_pdf_preview(pdf_path, page_num=0):
    """Extract preview image from PDF page"""
    try:
        pdf_document = fitz.open(pdf_path)
        page = pdf_document[page_num]
        
        # Render page as image
        mat = fitz.Matrix(1.0, 1.0)  # Scale factor
        pix = page.get_pixmap(matrix=mat)
        img_data = pix.tobytes("png")
        
        pdf_document.close()
        return img_data
    
    except Exception as e:
        logger.error(f"PDF preview extraction error: {str(e)}")
        return None